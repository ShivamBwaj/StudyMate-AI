# ppt_summarizer.py

import os
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()
api_key = os.getenv("GROQ_API_KEY")

client = OpenAI(
    api_key=api_key,
    base_url="https://api.groq.com/openai/v1"
)

def extract_text_from_pptx(pptx_path):
    """Extract all text from PowerPoint slides."""
    prs = Presentation(pptx_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text.strip()

def summarize_pptx(pptx_path):
    """Extract and summarize slide text using LLaMA-3."""
    text = extract_text_from_pptx(pptx_path)
    
    if not text:
        return "❌ No readable text found in the PPTX."

    if len(text) > 10000:
        text = text[:10000]  # Truncate if too long

    print("🧠 Summarizing PowerPoint...")

    response = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {"role": "system", "content": "You are a helpful assistant summarizing lecture slides."},
            {"role": "user", "content": text}
        ],
        temperature=0.5,
    )

    return response.choices[0].message.content
